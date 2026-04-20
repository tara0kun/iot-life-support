"""家族向け説明用パワーポイントを生成する。

実行: python scripts/build_proposal.py
出力: proposal.pptx (プロジェクトルート)
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt

OUT = Path(__file__).resolve().parent.parent / "proposal.pptx"

# 色定義
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0xE8, 0x8B, 0x3B)  # 温かいオレンジ
GREEN = RGBColor(0x4A, 0x7C, 0x59)
RED = RGBColor(0xB5, 0x3B, 0x3B)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xF7, 0xF3, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Cm(33.867)   # 16:9
prs.slide_height = Cm(19.05)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, text, left, top, width, height, *,
             font_size=18, bold=False, color=NAVY, align="left", font="Meiryo"):
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.2); tf.margin_right = Cm(0.2)
    tf.margin_top = Cm(0.1); tf.margin_bottom = Cm(0.1)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_title_bar(slide, title, num=None, total=None):
    # 上部のタイトル帯
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(2.2))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, title, Cm(1.2), Cm(0.4), SLIDE_W - Cm(4), Cm(1.5),
             font_size=28, bold=True, color=WHITE)
    if num is not None:
        add_text(slide, f"{num} / {total}", SLIDE_W - Cm(3), Cm(0.7),
                 Cm(2.5), Cm(0.8), font_size=14, color=WHITE, align="right")


def add_box(slide, text, left, top, width, height, *,
            fill=WHITE, border=NAVY, font_size=16, title=None,
            title_color=NAVY, body_color=GRAY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.06
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.4); tf.margin_right = Cm(0.4)
    tf.margin_top = Cm(0.3); tf.margin_bottom = Cm(0.3)
    first = True
    if title:
        p = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.name = "Meiryo"
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
        first = False
    for line in text.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.name = "Meiryo"
        run.font.size = Pt(font_size)
        run.font.color.rgb = body_color


# ============================================================
# Slide 1: 表紙
# ============================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)
# アクセント帯
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(13), SLIDE_W, Cm(0.3))
accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()

add_text(s, "おばあちゃんの", Cm(2), Cm(5), SLIDE_W - Cm(4), Cm(2),
         font_size=40, bold=True, color=WHITE)
add_text(s, "毎日を、そっと支える仕組み", Cm(2), Cm(7.2), SLIDE_W - Cm(4), Cm(2),
         font_size=40, bold=True, color=ACCENT)
add_text(s, "IoTを使った生活サポートシステムのご提案",
         Cm(2), Cm(10), SLIDE_W - Cm(4), Cm(1.2),
         font_size=20, color=WHITE)
add_text(s, "── 家族みんなで、おばあちゃんの尊厳と安心を両立するために ──",
         Cm(2), Cm(13.8), SLIDE_W - Cm(4), Cm(1),
         font_size=14, color=WHITE)
add_text(s, "2026年4月  |  家族会議用資料",
         Cm(2), Cm(16.5), SLIDE_W - Cm(4), Cm(1),
         font_size=13, color=WHITE)

TOTAL = 14  # 後で更新

# ============================================================
# Slide 2: 今、家で起きていること
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "今、おばあちゃんの生活で起きていること", 2, TOTAL)

add_text(s, "認知症の進行により、以下のような場面が増えています",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1),
         font_size=18, color=GRAY)

items = [
    ("🍚", "ご飯を何度も食べてしまう",
     "・朝食のあとすぐに、また炊飯器を開けてしまう\n・冷蔵庫から何度もおかずを出して食べる\n・お腹を壊す・栄養バランスが崩れる"),
    ("🚪", "訪問販売への対応",
     "・見覚えのない契約書が後から出てくる\n・来訪者のことを数分で忘れてしまう\n・同じ人が何度も来ていることに気づけない"),
    ("🚿", "入浴・洗髪のこと",
     "・お風呂に入っても頭を洗わない日がある\n・「洗った」と本人は言う\n・指摘すると怒ってしまう"),
]
y = Cm(4.2)
for emoji, head, body in items:
    add_box(s, body, Cm(1.5), y, SLIDE_W - Cm(3), Cm(3.8),
            fill=WHITE, border=ACCENT, title=f"{emoji}  {head}",
            title_color=NAVY, body_color=GRAY, font_size=15)
    y += Cm(4.2)

# ============================================================
# Slide 3: 家族が困っていること
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "家族のジレンマ", 3, TOTAL)

add_text(s, "指摘したくない、でも放っておけない",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=22, bold=True, color=NAVY)

# 左右の比較
left_x = Cm(1.5); right_x = Cm(17.5); boxw = Cm(14.5); boxh = Cm(10.5)
add_box(s,
        "「また炊いたの？」\n「さっき食べたでしょ」\n"
        "\n→ おばあちゃんは傷つき、怒り、\n   家族の関係がぎくしゃくする\n"
        "\n→ プライドを傷つけると、\n   さらに強く反発してしまう",
        left_x, Cm(4.5), boxw, boxh,
        fill=WHITE, border=RED, title="❌  直接指摘すると",
        title_color=RED, body_color=GRAY, font_size=16)

add_box(s,
        "24時間見張ることはできない\n"
        "\n目を離した数分で\n・炊飯器を何度も開ける\n・訪問販売と契約する\n・火を扱う\n"
        "\n→ 誰かが常に付いていないと不安",
        right_x, Cm(4.5), boxw, boxh,
        fill=WHITE, border=RED, title="❌  何も言わないと",
        title_color=RED, body_color=GRAY, font_size=16)

add_text(s, "この2つの間で、技術の力を借りられないか？",
         Cm(1.5), Cm(16.5), SLIDE_W - Cm(3), Cm(1.2),
         font_size=20, bold=True, color=ACCENT, align="center")

# ============================================================
# Slide 4: 提案するシステム概要
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "ご提案：IoTによる「そっと支える」仕組み", 4, TOTAL)

add_text(s,
         "おばあちゃんを「見張る」のではなく、",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1),
         font_size=22, color=NAVY)
add_text(s,
         "「自分の記録帳」をそっと差し出すシステムです",
         Cm(1.5), Cm(4), SLIDE_W - Cm(3), Cm(1),
         font_size=22, bold=True, color=ACCENT)

# 3つの柱
add_box(s,
        "リビングのタブレットに\n"
        "・さっき何時にご飯を食べた\n"
        "・冷蔵庫は朝から何回開けた\n"
        "などを、おばあちゃん自身が\n見られるようにする\n"
        "\n「あら、もう食べたんだ」と\n本人が自然に気づける",
        Cm(1.5), Cm(6), Cm(10), Cm(11),
        fill=WHITE, border=GREEN, title="①  自分で気づく",
        title_color=GREEN, body_color=GRAY, font_size=15)

add_box(s,
        "炊飯器やIHは\n食後しばらく「調子が悪い」\n状態になる\n"
        "\nおばあちゃんは\n「壊れてるのかしら」\nと思うだけ\n"
        "\n→ 責めない、気づかれない\n   物理的な安全弁",
        Cm(12), Cm(6), Cm(10), Cm(11),
        fill=WHITE, border=GREEN, title="②  機械のフリで止める",
        title_color=GREEN, body_color=GRAY, font_size=15)

add_box(s,
        "おばあちゃんが同じことを\n繰り返そうとしたら\n"
        "家族のLINEに通知\n"
        "\n家族が\n「あれ、おばあちゃん、\nコーヒーでも飲もうよ」と\n自然に声をかけられる",
        Cm(22.5), Cm(6), Cm(10), Cm(11),
        fill=WHITE, border=GREEN, title="③  家族にそっと連絡",
        title_color=GREEN, body_color=GRAY, font_size=15)

# ============================================================
# Slide 5: 設計の絶対原則
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "このシステムの「絶対に守ること」", 5, TOTAL)

add_text(s, "おばあちゃんの尊厳を何より大切にします",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=20, bold=True, color=NAVY)

# 2x2 のルール表
rules = [
    ("❌ やらないこと", "「もう食べたよ」と指摘する", RED),
    ("✅ 代わりに", "本人が見る記録帳として見せる", GREEN),
    ("❌ やらないこと", "「また炊いたの？」と責める", RED),
    ("✅ 代わりに", "機械の不調に見せて自然に止める", GREEN),
    ("❌ やらないこと", "監視システムと分からせる", RED),
    ("✅ 代わりに", "「自分の記録帳」として提示する", GREEN),
    ("❌ やらないこと", "家族が止めていると思わせる", RED),
    ("✅ 代わりに", "機械が壊れているだけ、と自然に", GREEN),
]
y = Cm(4.5); i = 0
for head, body, color in rules:
    col = i % 2
    row = i // 2
    x = Cm(1.5) + Cm(16) * col
    add_box(s, body, x, y + Cm(3.1) * row, Cm(15), Cm(2.8),
            fill=WHITE, border=color, title=head,
            title_color=color, body_color=GRAY, font_size=14)
    i += 1

# ============================================================
# Slide 6: 具体例：ご飯問題への対応
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "例：ご飯を何度も食べてしまう問題", 6, TOTAL)

add_text(s, "4つの段階で、少しずつ介入を強めていきます",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=18, color=GRAY)

stages = [
    ("第1段階", "気づかせる",
     "タブレットに「さっき 8:15 に朝ご飯を食べました」と表示。おばあちゃん自身が見て「あ、そうだった」と思い出せる"),
    ("第2段階", "機械のフリで止める",
     "1回目の食事後、炊飯器とIHは一時的に電源が入らない。「壊れたかしら」と思ってもらう"),
    ("第3段階", "家族にLINE通知",
     "2回目の食事行動を検知したら、母や祖父のLINEに通知。自然な声かけで介入"),
    ("第4段階", "おばあちゃん専用の食料",
     "冷蔵庫に最初から「おばあちゃんプレート」を用意。食べ過ぎになる食料を最初から置かない（運用面の工夫）"),
]
y = Cm(4.5)
for name, head, body in stages:
    add_box(s, body, Cm(1.5), y, SLIDE_W - Cm(3), Cm(2.8),
            fill=WHITE, border=NAVY, title=f"{name}  ●  {head}",
            title_color=NAVY, body_color=GRAY, font_size=14)
    y += Cm(3.1)

# ============================================================
# Slide 7: 家族の使い方（プライバシー設計）
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "家族はどう使う？", 7, TOTAL)

add_text(s, "おばあちゃんと家族で、見える情報が完全に分かれます",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=18, color=GRAY)

# 比較表（横並び）
add_box(s,
        "📍 リビングの壁掛けタブレット\n"
        "\n【見える情報】\n自分の行動記録だけ\n"
        "\n【できること】\n・記録を見る\n・「今から使います」と宣言\n"
        "\n【見えないもの】\n・家族の行動\n・家族の編集操作\n・システムの設定",
        Cm(1.5), Cm(4.5), Cm(15), Cm(12),
        fill=WHITE, border=ACCENT, title="👵 おばあちゃん",
        title_color=ACCENT, body_color=GRAY, font_size=14)

add_box(s,
        "📱 自分のスマホから\n"
        "\n【見える情報】\n家族全員の記録\n"
        "\n【できること】\n・誤った記録を修正\n・機器のロック解除\n・通知の確認\n"
        "\n【重要】\nこれらの操作は、\nおばあちゃんには\n絶対に見えません",
        Cm(17.5), Cm(4.5), Cm(15), Cm(12),
        fill=WHITE, border=GREEN, title="👨‍👩 家族（お母さん・おじいちゃん等）",
        title_color=GREEN, body_color=GRAY, font_size=14)

# ============================================================
# Slide 8: どうやって人を区別するか
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "誰が使っているかを、どう区別するか", 8, TOTAL)

add_text(s, "カメラによる顔認識 ＋ タブレットでの宣言",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=22, bold=True, color=NAVY)

add_text(s,
         "キッチンに小さなカメラを設置し、事前に登録した家族の顔を認識します。\n"
         "認識できない時は、タブレットで「誰が使うか」をボタンで選ぶ仕組みです。",
         Cm(1.5), Cm(4.3), SLIDE_W - Cm(3), Cm(2.2),
         font_size=16, color=GRAY)

scenarios = [
    ("おばあちゃんが冷蔵庫に近づく",
     "カメラが認識 → 直近の食事をチェック → 最近食べていたら「さっき召し上がりましたよ」と表示 → 本人が[使う]を押せば開く"),
    ("お母さんが冷蔵庫に近づく",
     "カメラが認識 → 無言で開く。おばあちゃんの画面には一切表示されない"),
    ("カメラが認識できない（暗い・横向き等）",
     "タブレットで「誰が使いますか？」を表示 → ボタンで選択 → 宣言した情報で進む"),
]
y = Cm(7); i = 0
colors = [ACCENT, GREEN, NAVY]
for head, body in scenarios:
    add_box(s, body, Cm(1.5), y, SLIDE_W - Cm(3), Cm(3),
            fill=WHITE, border=colors[i], title=head,
            title_color=colors[i], body_color=GRAY, font_size=14)
    y += Cm(3.3); i += 1

# ============================================================
# Slide 9: プライバシー・倫理への配慮
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "プライバシー・倫理への配慮", 9, TOTAL)

add_text(s, "心配な点に、事前に正直に向き合います",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1.2),
         font_size=20, bold=True, color=NAVY)

concerns = [
    ("Q. カメラの映像は保存される？",
     "A. 映像そのものは保存しません。『いつ・誰が・何をした』という記録だけをデータベースに残します。家族の判断で保存期間を決められます。"),
    ("Q. 外部に情報が漏れない？",
     "A. すべて自宅のラズパイ内に保存されます。インターネット経由で家族が見る場合も、VPN等で暗号化します。"),
    ("Q. おばあちゃんに黙ってやるのは倫理的に問題では？",
     "A. 重要な論点です。完全に秘密にするか、ある程度説明するかは家族で話し合いましょう。認知症ケアでは『記録帳』として緩やかに伝えるのが一般的です。"),
    ("Q. 祖父も映り込むのでは？",
     "A. その通りです。祖父にも事前に説明し、同意を得る必要があります。"),
]
y = Cm(4.3)
for q, a in concerns:
    add_box(s, a, Cm(1.5), y, SLIDE_W - Cm(3), Cm(3),
            fill=WHITE, border=NAVY, title=q,
            title_color=NAVY, body_color=GRAY, font_size=13)
    y += Cm(3.25)

# ============================================================
# Slide 10: 必要な機材と費用
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "必要な機材と費用", 10, TOTAL)

add_text(s, "段階的に導入できます。最初の投資は小さく。",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1),
         font_size=18, color=GRAY)

add_box(s,
        "・Tapo P115（スマートプラグ） × 1    約 2,000円\n"
        "・Tapo T110（開閉センサー） × 2     約 5,000円\n"
        "・Tapo H100（センサー用ハブ） × 1   約 3,500円\n"
        "\n合計  約 10,000円\n"
        "\n炊飯器だけで効果があるか確認",
        Cm(1.5), Cm(4.3), Cm(10.3), Cm(11.5),
        fill=WHITE, border=GREEN,
        title="第1段階：お試し",
        title_color=GREEN, body_color=GRAY, font_size=13)

add_box(s,
        "・第1段階のすべて\n"
        "・Tapo T110 追加 × 2（冷蔵庫等）\n"
        "・Tapo C210（カメラ） × 1\n"
        "・Tapo P115 追加 × 1（IH）\n"
        "・タブレット × 1（既存流用可）\n"
        "\n追加 約 30,000〜45,000円\n累計 約 40,000〜55,000円",
        Cm(12), Cm(4.3), Cm(10.3), Cm(11.5),
        fill=WHITE, border=NAVY,
        title="第2段階：食事問題に本格対応",
        title_color=NAVY, body_color=GRAY, font_size=13)

add_box(s,
        "・Tapo D230S1（玄関ドアホン） × 1\n"
        "・SwitchBot Lock × 1〜2\n"
        "・スマートスピーカー（任意）\n"
        "\n追加 約 30,000〜60,000円\n"
        "\n訪問者問題・冷蔵庫物理ロック\n入浴（ドライヤー）対応",
        Cm(22.5), Cm(4.3), Cm(10.3), Cm(11.5),
        fill=WHITE, border=ACCENT,
        title="第3段階：その他の問題も",
        title_color=ACCENT, body_color=GRAY, font_size=13)

add_text(s, "※ ラズベリーパイ本体・既存タブレットは手元のものを使用します",
         Cm(1.5), Cm(16.8), SLIDE_W - Cm(3), Cm(0.8),
         font_size=12, color=GRAY)

# ============================================================
# Slide 11: 導入スケジュール案
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "導入スケジュール案", 11, TOTAL)

add_text(s, "家族の同意をいただいてから、慎重に段階的に進めます",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1),
         font_size=18, color=GRAY)

steps = [
    ("① 今", "家族会議（この資料）",
     "この資料を元に、家族全員で議論・合意形成"),
    ("② 1〜2週間", "自宅で試験",
     "私（孫）の自宅で炊飯器・センサーを動かし、動作を確認"),
    ("③ 1ヶ月目", "実家に第1段階を導入",
     "炊飯器だけを対象に、記録機能のみ稼働。おばあちゃんの反応を観察"),
    ("④ 2〜3ヶ月目", "拡張",
     "効果があれば冷蔵庫・IH・カメラへ拡張。並行して家族UIを整備"),
    ("⑤ 3ヶ月目以降", "調整・運用",
     "数値を見ながら介入の強さを調整。他の問題（訪問者・入浴）にも取り組む"),
]
y = Cm(4.3)
for when, head, body in steps:
    # 左側にステップバッジ
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1.5), y, Cm(5), Cm(2.5))
    badge.adjustments[0] = 0.15
    badge.fill.solid(); badge.fill.fore_color.rgb = NAVY; badge.line.fill.background()
    tf = badge.text_frame; tf.margin_left = Cm(0.2); tf.margin_right = Cm(0.2)
    from pptx.enum.text import PP_ALIGN
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = when; r.font.name = "Meiryo"; r.font.size = Pt(14)
    r.font.bold = True; r.font.color.rgb = WHITE
    # 本文
    add_box(s, body, Cm(7), y, Cm(25.8), Cm(2.5),
            fill=WHITE, border=NAVY, title=head,
            title_color=NAVY, body_color=GRAY, font_size=13)
    y += Cm(2.8)

# ============================================================
# Slide 12: リスクと対策
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "想定されるリスクと対策", 12, TOTAL)

risks = [
    ("おばあちゃんがシステムに気づいて怒る",
     "設計原則を厳守。機械の不調・記録帳として緩やかに導入。うまくいかなければ撤去する柔軟性を持つ"),
    ("カメラが誤認識して家族の行動を祖母の記録にしてしまう",
     "家族用UIから簡単に修正可能。修正しても祖母には何も見えない"),
    ("機器が故障してロックが解除されない",
     "家族のスマホから緊急解除が可能。ロック機構が壊れた場合は『常時開』側に倒れる設計"),
    ("システムが原因でトラブルが起きる",
     "1ヶ月ごとに家族で振り返り。問題があれば一部機能または全体を停止する基準を事前に決めておく"),
    ("インターネットが切れて家族通知が届かない",
     "記録機能はローカルで動き続ける。通知はあとから届く。重要な機能はネット依存しない設計"),
]
y = Cm(3.3)
for risk, mitigation in risks:
    add_box(s, mitigation, Cm(1.5), y, SLIDE_W - Cm(3), Cm(2.7),
            fill=WHITE, border=RED, title=f"⚠  {risk}",
            title_color=RED, body_color=GRAY, font_size=13)
    y += Cm(2.95)

# ============================================================
# Slide 13: 家族に決めていただきたいこと
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_title_bar(s, "家族で決めていただきたいこと", 13, TOTAL)

add_text(s, "次回までに、以下について話し合ってください",
         Cm(1.5), Cm(2.8), SLIDE_W - Cm(3), Cm(1),
         font_size=18, color=NAVY)

decisions = [
    ("🎯 取り組む問題の優先順位",
     "食事／訪問者／入浴 の3つのうち、まずどれから始めるか"),
    ("💴 初期投資の予算",
     "第1段階（1万円）だけで様子見か、いきなり第2段階（4〜5万円）まで進めるか"),
    ("🤝 おばあちゃんへの説明",
     "完全に秘密にするか／「記録帳」として緩やかに伝えるか／全部説明するか"),
    ("👥 家族の役割分担",
     "誰が日常の確認をする？ 誰が通知に対応する？ 誰が記録を修正する？"),
    ("📹 カメラへの同意",
     "祖父の同意はどう取る？ 映像保存の期間は？"),
    ("🛑 中断の基準",
     "どういう状況になったら、このシステムを止める／見直すか"),
]
y = Cm(4.3); i = 0
for head, body in decisions:
    col = i % 2; row = i // 2
    x = Cm(1.5) + Cm(16) * col
    add_box(s, body, x, y + Cm(4) * row, Cm(15), Cm(3.7),
            fill=WHITE, border=ACCENT, title=head,
            title_color=ACCENT, body_color=GRAY, font_size=13)
    i += 1

# ============================================================
# Slide 14: 最後に
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s, NAVY)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(13), SLIDE_W, Cm(0.3))
accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background()

add_text(s, "大切なのは、", Cm(2), Cm(4), SLIDE_W - Cm(4), Cm(1.5),
         font_size=32, color=WHITE)
add_text(s, "技術で置き換えることではなく、", Cm(2), Cm(6), SLIDE_W - Cm(4), Cm(1.5),
         font_size=32, bold=True, color=WHITE)
add_text(s, "家族のやさしさを、少しだけ楽にすること。", Cm(2), Cm(8.2), SLIDE_W - Cm(4), Cm(1.5),
         font_size=32, bold=True, color=ACCENT)

add_text(s,
         "おばあちゃんが、\nできるだけ長く、自分らしく暮らせるように。\n"
         "家族が、できるだけ疲弊せずに支えられるように。",
         Cm(2), Cm(14), SLIDE_W - Cm(4), Cm(3),
         font_size=16, color=WHITE)

prs.save(OUT)
print(f"作成完了: {OUT}")
print(f"サイズ: {OUT.stat().st_size // 1024} KB / {len(prs.slides)} スライド")
